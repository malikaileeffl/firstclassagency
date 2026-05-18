"""Assemble final FC April Week 5.pptx from the 3 rendered slide PNGs."""
import os
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor

ASSETS = "/sessions/adoring-dazzling-babbage/mnt/outputs/assets"
OUT_PPT = "/sessions/adoring-dazzling-babbage/mnt/outputs/FC April Week 5.pptx"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

blank = prs.slide_layouts[6]   # blank layout

for i, name in enumerate(["slide1.png", "slide2.png", "slide3.png"]):
    slide = prs.slides.add_slide(blank)
    img_path = os.path.join(ASSETS, name)
    pic = slide.shapes.add_picture(img_path, 0, 0,
                                   width=prs.slide_width, height=prs.slide_height)

prs.save(OUT_PPT)
print("Saved", OUT_PPT)
